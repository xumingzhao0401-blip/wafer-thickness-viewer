# -*- coding: utf-8 -*-
"""
Wafer Thickness Viewer (Streamlit) - v3 (PPT layout + radius column + one-click download)

Changes requested:
1) PPT: "Generate & Download" should download directly (use st.download_button with prebuilt bytes).
2) PPT layout: left two images stacked, each fixed 8cm x 8cm; middle raw table; right stats table; no overlap.
3) Tables (web & PPT): show columns as X, Y, Radius, Thickness (Radius computed from X,Y). X/Y/Radius keep 2 decimals.
"""

from __future__ import annotations

from typing import List, Optional, Tuple
from datetime import datetime
import io
import math

import numpy as np
import pandas as pd
import plotly.graph_objects as go
import plotly.io as pio
import streamlit as st

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# -----------------------------
# Page config
# -----------------------------
st.set_page_config(page_title="晶圆厚度 3D 可视化", page_icon="🟢", layout="wide")


# ============================================================
# Dataset Manager
# ============================================================

def _init_dataset_store():
    if "datasets" not in st.session_state:
        st.session_state.datasets = {}
    if "active_ds" not in st.session_state:
        st.session_state.active_ds = None


def _unique_name(base: str) -> str:
    base = (base or "Wafer").strip()
    name = base
    i = 2
    while name in st.session_state.datasets:
        name = f"{base} ({i})"
        i += 1
    return name


def register_dataset(
    name: str,
    df: pd.DataFrame,
    wafer_inch: float,
    cmap: str = "viridis",
    spec_upper: Optional[float] = None,
    spec_lower: Optional[float] = None,
):
    _init_dataset_store()
    name = _unique_name(name)
    st.session_state.datasets[name] = {
        "df": df.copy(),  # keep x,y,thickness in μm
        "wafer_inch": float(wafer_inch),
        "cmap": cmap,
        "spec_upper": spec_upper,
        "spec_lower": spec_lower,
        "created_at": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
    }
    st.session_state.active_ds = name


def get_active_dataset() -> Tuple[Optional[str], Optional[dict]]:
    _init_dataset_store()
    name = st.session_state.active_ds
    if not name:
        return None, None
    return name, st.session_state.datasets.get(name)


def rename_dataset(old_name: str, new_name: str):
    _init_dataset_store()
    if old_name not in st.session_state.datasets:
        return
    new_name = (new_name or "").strip()
    if not new_name or new_name == old_name:
        return
    new_name = _unique_name(new_name)
    st.session_state.datasets[new_name] = st.session_state.datasets.pop(old_name)
    st.session_state.active_ds = new_name


def delete_dataset(name: str):
    _init_dataset_store()
    st.session_state.datasets.pop(name, None)
    st.session_state.active_ds = next(iter(st.session_state.datasets), None)


# ============================================================
# Units
# ============================================================

def convert_to_um(values: pd.Series, unit: str) -> pd.Series:
    """Convert thickness to μm. 1 μm = 10000 Å."""
    unit = unit.strip()
    if unit == "μm":
        return values
    if unit in ["Å", "A", "Angstrom", "angstrom", "Ångström"]:
        return values / 10000.0
    return values


# ============================================================
# Utils / Core math (Cached)
# ============================================================

def build_grids(radius_mm: float, grid_res: int) -> Tuple[np.ndarray, np.ndarray]:
    x_lin = np.linspace(-radius_mm, radius_mm, grid_res)
    y_lin = np.linspace(-radius_mm, radius_mm, grid_res)
    return np.meshgrid(x_lin, y_lin)


@st.cache_data
def cached_idw_interpolation(
    xs: np.ndarray,
    ys: np.ndarray,
    ts: np.ndarray,
    radius_mm: float,
    grid_res: int,
    power: float = 2.0,
    eps: float = 1e-6,
) -> Tuple[Optional[np.ndarray], Optional[np.ndarray], Optional[np.ndarray]]:
    if len(xs) < 3:
        return None, None, None

    grid_x, grid_y = build_grids(radius_mm, grid_res)

    gx = grid_x[..., np.newaxis]
    gy = grid_y[..., np.newaxis]
    px = xs[np.newaxis, np.newaxis, :]
    py = ys[np.newaxis, np.newaxis, :]

    dx = gx - px
    dy = gy - py
    dist = np.sqrt(dx**2 + dy**2) + eps

    weights = 1.0 / (dist**power)
    w_sum = np.sum(weights, axis=2)
    t_weighted_sum = np.sum(weights * ts[np.newaxis, np.newaxis, :], axis=2)

    with np.errstate(divide="ignore", invalid="ignore"):
        grid_t = t_weighted_sum / w_sum

    mask = (grid_x**2 + grid_y**2) <= (radius_mm**2)
    grid_t = np.where(mask, grid_t, np.nan)

    return grid_x, grid_y, grid_t


@st.cache_data
def compute_statistics(ts: np.ndarray) -> pd.DataFrame:
    """
    Unit effect:
    - mean/std/min/max/p-p scale with unit
    - percentage metrics (CoV, WIWNU...) are unit-invariant
    """
    ts = np.asarray(ts, dtype=float)
    ts = ts[~np.isnan(ts)]

    n = int(ts.size)
    if n == 0:
        return pd.DataFrame(columns=["特征名", "公式", "计算结果", "单位"])

    mean = float(np.mean(ts))
    std = float(np.std(ts, ddof=1)) if n > 1 else 0.0
    min_v = float(np.min(ts))
    max_v = float(np.max(ts))
    p2p = max_v - min_v

    safe_mean = mean if mean != 0 else 1e-9
    safe_sum = (max_v + min_v) if (max_v + min_v) != 0 else 1e-9

    cov = (std / safe_mean) * 100.0
    wiw_nu = (p2p / safe_sum) * 100.0
    range_mean = (p2p / safe_mean) * 100.0
    max_dev = (float(np.max(np.abs(ts - mean))) / safe_mean) * 100.0
    u3 = (3.0 * std / safe_mean) * 100.0
    u6 = (6.0 * std / safe_mean) * 100.0

    rows = [
        ("点数 N", "N", n, ""),
        ("平均值 (Mean)", "μ = (1/N)·Σxi", mean, "μm"),
        ("标准差 (Std, 1σ)", "σ = sqrt( Σ(xi-μ)^2 / (N-1) )", std, "μm"),
        ("变异系数 (CoV)", "CoV = σ/μ × 100%", cov, "%"),
        ("最小值 (Min)", "min(xi)", min_v, "μm"),
        ("最大值 (Max)", "max(xi)", max_v, "μm"),
        ("峰-峰值 (Peak-to-Peak)", "P-P = max - min", p2p, "μm"),
        ("WIWNU（常用均匀度）", "(max-min)/(max+min) × 100%", wiw_nu, "%"),
        ("Range/Mean", "(max-min)/μ × 100%", range_mean, "%"),
        ("最大偏差 (Max Dev.)", "max(|xi-μ|)/μ × 100%", max_dev, "%"),
        ("3σ 均匀度", "3σ/μ × 100%", u3, "%"),
        ("6σ 均匀度", "6σ/μ × 100%", u6, "%"),
    ]

    df = pd.DataFrame(rows, columns=["特征名", "公式", "计算结果", "单位"])
    df["计算结果"] = df["计算结果"].apply(
        lambda v: f"{v:.4f}" if isinstance(v, (float, np.floating)) else str(v)
    )
    return df


@st.cache_data
def load_csv(uploaded_file) -> Optional[pd.DataFrame]:
    try:
        uploaded_file.seek(0)
        df = pd.read_csv(uploaded_file)
    except Exception as e:
        st.error(f"读取 CSV 失败：{e}")
        return None

    lower_map = {c.lower().strip(): c for c in df.columns}
    required = ["x", "y", "thickness"]
    missing = [k for k in required if k not in lower_map]
    if missing:
        st.error(f"CSV 缺少列：{missing}")
        return None

    df2 = df[[lower_map["x"], lower_map["y"], lower_map["thickness"]]].copy()
    df2.columns = ["x", "y", "thickness"]
    for c in required:
        df2[c] = pd.to_numeric(df2[c], errors="coerce")

    df2 = df2.dropna(subset=["x", "y", "thickness"]).reset_index(drop=True)
    if df2.empty:
        st.error("CSV 解析后没有有效数据。")
        return None
    return df2


def wafer_radius_mm(wafer_inch: float) -> float:
    return float(wafer_inch) * 25.4 / 2.0


def circle_boundary_trace(radius_mm: float, n: int = 361) -> go.Scatter:
    theta = np.linspace(0, 2 * np.pi, n)
    x = radius_mm * np.cos(theta)
    y = radius_mm * np.sin(theta)
    return go.Scatter(x=x, y=y, mode="lines", showlegend=False, line=dict(width=3, color="black"))


def with_radius_columns(df: pd.DataFrame) -> pd.DataFrame:
    """Return a display-ready DF with X,Y,Radius,Thickness; keep rounding for display."""
    d = df.copy()
    d["radius"] = np.sqrt(d["x"].astype(float) ** 2 + d["y"].astype(float) ** 2)
    # display rounding
    d["x"] = d["x"].astype(float).round(2)
    d["y"] = d["y"].astype(float).round(2)
    d["radius"] = d["radius"].astype(float).round(2)
    d["thickness"] = d["thickness"].astype(float).round(4)
    return d[["x", "y", "radius", "thickness"]]


# ============================================================
# Drawing Functions
# ============================================================

def make_top_view_heatmap(
    df: pd.DataFrame,
    radius_mm: float,
    cmap: str,
    grid_res: int,
    show_labels: bool = True,
    view_center: Tuple[float, float] = (0.0, 0.0),
    view_zoom: float = 1.0,
) -> go.Figure:
    xs = df["x"].to_numpy(dtype=float)
    ys = df["y"].to_numpy(dtype=float)
    ts = df["thickness"].to_numpy(dtype=float)

    grid_x, grid_y, grid_t = cached_idw_interpolation(xs, ys, ts, radius_mm, grid_res)

    fig = go.Figure()
    if grid_t is not None:
        fig.add_trace(go.Heatmap(
            z=grid_t,
            x=grid_x[0, :],
            y=grid_y[:, 0],
            colorscale=cmap,
            colorbar=dict(title="Thickness (μm)", thickness=18, len=0.85),
            hovertemplate="X=%{x:.2f}<br>Y=%{y:.2f}<br>T=%{z:.4f} μm<extra></extra>",
        ))

    fig.add_trace(circle_boundary_trace(radius_mm))

    fig.add_trace(go.Scatter(
        x=xs, y=ys, mode="markers",
        marker=dict(size=8, color="black", line=dict(width=1, color="white")),
        showlegend=False,
        hovertemplate="X=%{x:.2f}<br>Y=%{y:.2f}<br>T=%{text} μm<extra></extra>",
        text=[f"{v:.4f}" for v in ts],
    ))

    if show_labels:
        fig.add_trace(go.Scatter(
            x=xs, y=ys, mode="text",
            text=[f"{v:.4f}" for v in ts],
            textposition="middle center",
            showlegend=False,
            hoverinfo="skip",
            textfont=dict(color="black", size=10, family="Arial")
        ))

    cx, cy = view_center
    zoom = max(0.2, float(view_zoom))
    half = radius_mm * 1.05 / zoom
    xr = [cx - half, cx + half]
    yr = [cy - half, cy + half]

    fig.update_layout(
        title="Top View Heatmap (μm)",
        margin=dict(l=20, r=20, t=60, b=20),
        xaxis=dict(title="X (mm)", scaleanchor="y", scaleratio=1, range=xr),
        yaxis=dict(title="Y (mm)", range=yr),
        height=800,
    )
    return fig


def make_3d_surface(
    df: pd.DataFrame,
    radius_mm: float,
    cmap: str,
    grid_res: int,
    spec_upper: Optional[float],
    spec_lower: Optional[float],
    camera_eye: Optional[Tuple[float, float, float]] = None,
    z_scale: float = 1.0,
    z_aspect: float = 2.5,
) -> go.Figure:
    xs = df["x"].to_numpy(dtype=float)
    ys = df["y"].to_numpy(dtype=float)
    ts = df["thickness"].to_numpy(dtype=float)

    grid_x, grid_y, grid_t = cached_idw_interpolation(xs, ys, ts, radius_mm, grid_res)

    grid_z = None if grid_t is None else (grid_t * float(z_scale))
    pts_z = ts * float(z_scale)

    fig = go.Figure()

    if grid_t is not None:
        fig.add_trace(go.Surface(
            x=grid_x, y=grid_y, z=grid_z,
            surfacecolor=grid_t,
            colorscale=cmap,
            colorbar=dict(title="Thickness (μm)", thickness=18, len=0.85),
            customdata=grid_t,
            hovertemplate="X=%{x:.2f}<br>Y=%{y:.2f}<br>T=%{customdata:.4f} μm<extra></extra>",
        ))

        mask = ~np.isnan(grid_t)
        if spec_upper is not None and np.isfinite(spec_upper):
            z_up = np.where(mask, float(spec_upper) * float(z_scale), np.nan)
            fig.add_trace(go.Surface(x=grid_x, y=grid_y, z=z_up, opacity=0.25, showscale=False,
                                     colorscale=[[0, "red"], [1, "red"]], hoverinfo="skip", name="USL"))
        if spec_lower is not None and np.isfinite(spec_lower):
            z_lo = np.where(mask, float(spec_lower) * float(z_scale), np.nan)
            fig.add_trace(go.Surface(x=grid_x, y=grid_y, z=z_lo, opacity=0.25, showscale=False,
                                     colorscale=[[0, "red"], [1, "red"]], hoverinfo="skip", name="LSL"))

    fig.add_trace(go.Scatter3d(
        x=xs, y=ys, z=pts_z,
        mode="markers",
        marker=dict(size=4, color="black"),
        name="Points",
        customdata=ts,
        hovertemplate="X=%{x:.2f}<br>Y=%{y:.2f}<br>T=%{customdata:.4f} μm<extra></extra>",
    ))

    scene = dict(
        xaxis_title="X (mm)",
        yaxis_title="Y (mm)",
        zaxis_title=f"Thickness × {z_scale:g} (display)",
        aspectmode="manual",
        aspectratio=dict(x=1, y=1, z=float(z_aspect)),
        domain=dict(x=[0.0, 1.0], y=[0.0, 1.0]),
    )
    if camera_eye is not None:
        scene["camera"] = dict(eye=dict(x=camera_eye[0], y=camera_eye[1], z=camera_eye[2]))

    fig.update_layout(title="3D Distribution (μm, Z scaled)", scene=scene, margin=dict(l=0, r=0, t=50, b=0), height=900)
    return fig


# ============================================================
# Preview Plot (point_id)
# ============================================================

def plot_pattern_preview(df: pd.DataFrame, radius_mm: float) -> go.Figure:
    fig = go.Figure()
    fig.add_trace(circle_boundary_trace(radius_mm))
    labels = df["point_id"].astype(int).astype(str).tolist() if "point_id" in df.columns else [str(i) for i in df.index]
    fig.add_trace(go.Scatter(
        x=df["x"], y=df["y"],
        mode="markers+text",
        marker=dict(size=12, color="red"),
        text=labels,
        textposition="top center",
        textfont=dict(size=14, color="red", family="Arial Black"),
        showlegend=False,
    ))
    fig.update_layout(
        title="点位分布预览 (Preview)",
        width=800,
        height=700,
        margin=dict(l=20, r=20, t=40, b=20),
        xaxis=dict(range=[-radius_mm * 1.1, radius_mm * 1.1], scaleanchor="y", scaleratio=1, zeroline=True, showgrid=True),
        yaxis=dict(range=[-radius_mm * 1.1, radius_mm * 1.1], zeroline=True, showgrid=True),
        hovermode="closest",
    )
    return fig


# ============================================================
# Generator patterns
# ============================================================

def generate_pattern_coords(pattern_type, radius_mm, edge_exclude_mm, **kwargs) -> pd.DataFrame:
    effective_r = max(0.0, float(radius_mm) - float(edge_exclude_mm))
    points: List[Tuple[float, float]] = []
    point_id: Optional[List[int]] = None

    if pattern_type == "十字交叉 (Cross)":
        n_per_arm = int(kwargs.get("points_per_arm", 3))
        points.append((0.0, 0.0))
        if effective_r > 0:
            radii = np.linspace(0, effective_r, n_per_arm + 1)[1:]
            for r in radii:
                points.extend([(r, 0.0), (-r, 0.0), (0.0, r), (0.0, -r)])

    elif pattern_type == "同心圆 (Concentric)":
        n_rings = int(kwargs.get("n_rings", 3))
        pts_per_ring = int(kwargs.get("pts_per_ring", 8))
        points.append((0.0, 0.0))
        if effective_r > 0:
            radii = np.linspace(0, effective_r, n_rings + 1)[1:]
            for r in radii:
                angles = np.linspace(0, 2 * np.pi, pts_per_ring, endpoint=False)
                for ang in angles:
                    points.append((r * np.cos(ang), r * np.sin(ang)))

    elif pattern_type == "均匀网格 (Grid)":
        step = float(kwargs.get("grid_step", 30.0))
        xs = np.arange(0, effective_r + 0.1, step)
        xs = np.concatenate((-xs[:0:-1], xs))
        ys = xs.copy()
        for x in xs:
            for y in ys:
                if (x**2 + y**2) <= effective_r**2:
                    points.append((x, y))

    elif pattern_type == "FAB 25点模板 (Fab25)":
        s = effective_r / 3.0 if effective_r > 0 else 0.0
        layout = [
            (4,   0*s,  3*s),
            (5,  -1*s,  2*s), (3, 0*s, 2*s), (6,  1*s, 2*s),
            (10, -2*s,  1*s), (9, -1*s, 1*s), (2, 0*s, 1*s), (8,  1*s, 1*s), (7,  2*s, 1*s),
            (11, -3*s,  0*s), (12, -2*s, 0*s), (13, -1*s, 0*s), (1, 0*s, 0*s), (14, 1*s, 0*s), (15, 2*s, 0*s), (16, 3*s, 0*s),
            (21, -2*s, -1*s), (20, -1*s, -1*s), (19, 0*s, -1*s), (18, 1*s, -1*s), (17, 2*s, -1*s),
            (22, -1*s, -2*s), (23, 0*s, -2*s), (24, 1*s, -2*s),
            (25,  0*s, -3*s),
        ]
        layout_sorted = sorted(layout, key=lambda t: t[0])
        point_id = [pid for pid, _, _ in layout_sorted]
        points = [(x, y) for _, x, y in layout_sorted]

    df = pd.DataFrame(points, columns=["x", "y"])
    if point_id is not None:
        df.insert(0, "point_id", point_id)
    # add radius for display (mm)
    df["radius"] = np.sqrt(df["x"]**2 + df["y"]**2)
    df["thickness"] = np.nan
    return df


# ============================================================
# PPT Export Helpers
# ============================================================

def _cm(v: float) -> float:
    return v / 2.54  # cm to inches


def _fig_to_png_bytes(fig: go.Figure, scale: float = 2.0) -> bytes:
    return pio.to_image(fig, format="png", engine="kaleido", scale=scale)


def _set_run_font(run, size_pt: int = 12, bold: bool = False, color: Optional[Tuple[int, int, int]] = None):
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.name = "Microsoft YaHei"
    if color is not None:
        run.font.color.rgb = RGBColor(color[0], color[1], color[2])


def _add_table(slide, df: pd.DataFrame, x, y, w, h, font_size=10, header_fill=(240, 240, 240)):
    rows, cols = df.shape[0] + 1, df.shape[1]
    table_shape = slide.shapes.add_table(rows, cols, x, y, w, h)
    table = table_shape.table

    # header
    for j, col in enumerate(df.columns):
        cell = table.cell(0, j)
        cell.text = str(col)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(*header_fill)
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                _set_run_font(r, size_pt=font_size, bold=True)

    # body
    for i in range(df.shape[0]):
        for j in range(cols):
            cell = table.cell(i + 1, j)
            cell.text = str(df.iat[i, j])
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    _set_run_font(r, size_pt=font_size, bold=False)

    return table_shape


def build_ppt_one_slide(
    title: str,
    subtitle: str,
    raw_df: pd.DataFrame,
    stats_df: pd.DataFrame,
    top_png: bytes,
    surf_png: bytes,
) -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    m = Inches(0.35)
    gap = Inches(0.22)

    # Title
    title_box = slide.shapes.add_textbox(m, Inches(0.15), prs.slide_width - 2*m, Inches(0.5))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    _set_run_font(run, size_pt=24, bold=True, color=(25, 25, 25))

    sub_box = slide.shapes.add_textbox(m, Inches(0.62), prs.slide_width - 2*m, Inches(0.32))
    tf2 = sub_box.text_frame
    tf2.clear()
    p2 = tf2.paragraphs[0]
    r2 = p2.add_run()
    r2.text = subtitle
    _set_run_font(r2, size_pt=11, bold=False, color=(90, 90, 90))

    content_top = Inches(1.0)

    # Left images: fixed 8cm x 8cm
    img_size = Inches(_cm(8.0))
    img_gap = Inches(0.10)

    left_x = m
    img1_y = content_top
    img2_y = img1_y + img_size + img_gap

    slide.shapes.add_picture(io.BytesIO(top_png), left_x, img1_y, width=img_size, height=img_size)
    slide.shapes.add_picture(io.BytesIO(surf_png), left_x, img2_y, width=img_size, height=img_size)

    # Middle raw table, Right stats table
    usable_w = prs.slide_width - 2*m
    left_zone_w = img_size
    right_zone_w = Inches(3.55)  # stats
    mid_zone_w = usable_w - left_zone_w - right_zone_w - 2*gap

    mid_x = left_x + left_zone_w + gap
    right_x = mid_x + mid_zone_w + gap

    table_h = img_size * 2 + img_gap  # align with stacked images

    # Labels
    raw_label = slide.shapes.add_textbox(mid_x, content_top - Inches(0.22), mid_zone_w, Inches(0.22))
    rtf = raw_label.text_frame
    rtf.clear()
    pr = rtf.paragraphs[0]
    rr = pr.add_run()
    rr.text = "Raw Data"
    _set_run_font(rr, size_pt=11, bold=True, color=(25, 25, 25))

    stats_label = slide.shapes.add_textbox(right_x, content_top - Inches(0.22), right_zone_w, Inches(0.22))
    stf = stats_label.text_frame
    stf.clear()
    ps = stf.paragraphs[0]
    rs = ps.add_run()
    rs.text = "Statistics"
    _set_run_font(rs, size_pt=11, bold=True, color=(25, 25, 25))

    # Fit rows
    max_rows = 28
    raw_show = raw_df.copy()
    if len(raw_show) > max_rows:
        raw_show = raw_show.head(max_rows).copy()
        # add a small note in subtitle to avoid clutter

    _add_table(slide, raw_show, mid_x, content_top, mid_zone_w, table_h, font_size=9)
    _add_table(slide, stats_df, right_x, content_top, right_zone_w, table_h, font_size=9)

    bio = io.BytesIO()
    prs.save(bio)
    return bio.getvalue()


@st.cache_data(show_spinner=False)
def build_ppt_cached(
    title: str,
    subtitle: str,
    raw_csv: str,
    stats_csv: str,
    top_json: str,
    surf_json: str,
) -> bytes:
    raw_df = pd.read_csv(io.StringIO(raw_csv))
    stats_df = pd.read_csv(io.StringIO(stats_csv))

    top_fig = pio.from_json(top_json)
    surf_fig = pio.from_json(surf_json)

    top_png = _fig_to_png_bytes(top_fig, scale=2.0)
    surf_png = _fig_to_png_bytes(surf_fig, scale=2.0)

    return build_ppt_one_slide(title, subtitle, raw_df, stats_df, top_png, surf_png)


# ============================================================
# Generator UI
# ============================================================

def blind_mode_ui():
    _init_dataset_store()

    st.markdown("### 🛠️ 坐标生成器 (Generator Mode)")
    st.caption("选择测量图案自动生成坐标。无数据点留空即可。")

    col_ctrl, col_data = st.columns([1, 2], gap="large")

    with col_ctrl:
        st.subheader("1. 晶圆设置")
        wafer_inch = st.selectbox("尺寸 (Inch)", [4, 6, 8, 12], index=2, key="gen_inch")
        radius_mm = wafer_radius_mm(wafer_inch)
        ee_mm = st.number_input("Edge Exclusion (mm)", value=3.0, min_value=0.0, key="gen_ee")

        st.subheader("2. 图案选择")
        pat_type = st.radio(
            "生成方式",
            ["FAB 25点模板 (Fab25)", "十字交叉 (Cross)", "同心圆 (Concentric)", "均匀网格 (Grid)"],
            key="gen_pat",
        )

        gen_params = {}
        if pat_type == "十字交叉 (Cross)":
            gen_params["points_per_arm"] = st.slider("每臂点数", 1, 15, 3, key="gen_arm")
        elif pat_type == "同心圆 (Concentric)":
            gen_params["n_rings"] = st.slider("圈数", 1, 10, 3, key="gen_rings")
            gen_params["pts_per_ring"] = st.slider("每圈点数", 4, 32, 8, step=4, key="gen_ppr")
        elif pat_type == "均匀网格 (Grid)":
            gen_params["grid_step"] = st.number_input("网格间距 (mm)", value=30.0, min_value=5.0, key="gen_step")
        else:
            st.info("该模板点数与编号严格对齐 FAB 示意图（point_id=1~25）。")

        st.subheader("3. 数据集命名")
        default_name = f"GEN {pat_type.split(' ')[0]} {datetime.now().strftime('%H%M%S')}"
        st.text_input("保存为数据集名称", value=default_name, key="gen_ds_name")

        if st.button("🔄 生成坐标表", type="primary", use_container_width=True, key="gen_build"):
            df = generate_pattern_coords(pat_type, radius_mm, ee_mm, **gen_params)
            st.session_state.gen_df = df
            st.rerun()

    with col_data:
        st.subheader("4. 数据录入")
        if "gen_df" not in st.session_state:
            st.info("👈 请先生成坐标表")
            return

        df_curr = st.session_state.gen_df.copy()
        # keep radius updated
        df_curr["radius"] = np.sqrt(df_curr["x"]**2 + df_curr["y"]**2)
        st.session_state.gen_df = df_curr

        with st.expander("👁️ 点位预览（显示 point_id）", expanded=True):
            preview_fig = plot_pattern_preview(df_curr, radius_mm)
            st.plotly_chart(preview_fig, use_container_width=True)

        st.markdown(f"**点数：{len(df_curr)}**。厚度单位：μm（留空表示无数据）")

        column_config = {
            "x": st.column_config.NumberColumn("X (mm)", format="%.2f", disabled=True),
            "y": st.column_config.NumberColumn("Y (mm)", format="%.2f", disabled=True),
            "radius": st.column_config.NumberColumn("Radius (mm)", format="%.2f", disabled=True),
            "thickness": st.column_config.NumberColumn("Thickness (μm)", format="%.4f", required=False),
        }
        if "point_id" in df_curr.columns:
            column_config["point_id"] = st.column_config.NumberColumn("Point ID", disabled=True)

        # order columns nicely
        show_cols = ["x", "y", "radius", "thickness"]
        if "point_id" in df_curr.columns:
            show_cols = ["point_id"] + show_cols

        edited_df = st.data_editor(
            df_curr[show_cols],
            column_config=column_config,
            use_container_width=True,
            height=340,
            key="data_editor_gen",
        )

        if st.button("✅ 保存为数据集", type="primary", use_container_width=True, key="gen_save"):
            final_df = edited_df.dropna(subset=["thickness"]).copy()
            final_df["thickness"] = pd.to_numeric(final_df["thickness"], errors="coerce")
            final_df = final_df.dropna(subset=["thickness"])
            if final_df.empty:
                st.error("所有点厚度均为空！")
            else:
                # store only x,y,thickness
                register_dataset(
                    name=st.session_state.get("gen_ds_name", "Generator"),
                    df=final_df[["x", "y", "thickness"]].copy(),
                    wafer_inch=wafer_inch,
                )
                st.success(f"已保存为数据集：{st.session_state.active_ds}（有效点 {len(final_df)}）")
                st.rerun()

        csv = edited_df.to_csv(index=False).encode("utf-8")
        st.download_button("📥 下载 CSV 模板", csv, "template.csv", "text/csv", use_container_width=True, key="gen_dl")


# ============================================================
# Normal Mode UI + PPT Export
# ============================================================

def normal_mode_ui():
    _init_dataset_store()

    st.markdown("### 普通模式")
    st.info("坐标单位 mm；厚度内部统一 μm。表格显示：X、Y、Radius、Thickness。", icon="ℹ️")

    left, right = st.columns([1, 2], gap="large")

    with left:
        st.markdown("#### 数据集管理")
        names = list(st.session_state.datasets.keys())
        if names:
            default_idx = names.index(st.session_state.active_ds) if st.session_state.active_ds in names else 0
            chosen = st.selectbox("选择要显示的数据集", names, index=default_idx, key="ds_select")
            st.session_state.active_ds = chosen
            ds_name, ds = get_active_dataset()
        else:
            ds_name, ds = None, None
            st.info("还没有数据集：可导入 CSV 或在生成器中保存。", icon="🧊")

        if ds_name and ds:
            c1, c2 = st.columns([1, 1])
            with c1:
                new_name = st.text_input("重命名当前数据集", value=ds_name, key="ds_rename_txt")
                if st.button("✏️ 应用重命名", use_container_width=True, key="ds_rename_btn"):
                    rename_dataset(ds_name, new_name)
                    st.rerun()
            with c2:
                st.caption(f"创建时间：{ds.get('created_at','-')}")
                if st.button("🗑️ 删除当前数据集", use_container_width=True, key="ds_del_btn"):
                    delete_dataset(ds_name)
                    st.rerun()

        st.markdown("---")
        st.markdown("#### 导入 CSV 作为新数据集")
        uploaded = st.file_uploader("选择 CSV 文件（需列 x,y,thickness）", type=["csv"], key="normal_uploader")
        unit = st.radio("Thickness 单位", ["μm", "Å"], horizontal=True, key="import_unit")
        import_inch = st.selectbox("导入时晶圆尺寸（英寸）", [4, 6, 8, 12], index=2, key="import_wafer")
        default_import_name = f"{uploaded.name} {datetime.now().strftime('%H%M%S')}" if uploaded is not None else ""
        import_name = st.text_input("新数据集名称", value=default_import_name, key="import_name")

        if st.button("➕ 导入为新数据集", type="primary", use_container_width=True, disabled=(uploaded is None), key="import_btn"):
            df_new = load_csv(uploaded)
            if df_new is not None:
                df_new = df_new.copy()
                df_new["thickness"] = convert_to_um(df_new["thickness"], unit)
                register_dataset(name=import_name or uploaded.name, df=df_new[["x", "y", "thickness"]].copy(), wafer_inch=import_inch)
                st.success(f"已导入数据集：{st.session_state.active_ds}（thickness 已统一为 μm）")
                st.rerun()

        st.markdown("---")
        st.markdown("#### 绘图参数（作用于当前数据集）")
        if ds is None:
            st.info("请先选择/创建一个数据集。")
            return

        wafer_inch = st.selectbox("晶圆尺寸（英寸）", [4, 6, 8, 12],
                                 index=[4, 6, 8, 12].index(int(ds.get("wafer_inch", 8))), key="normal_wafer")
        radius_mm = wafer_radius_mm(wafer_inch)

        cmaps = ["viridis", "plasma", "inferno", "magma", "coolwarm", "cividis"]
        cmap = st.selectbox("色图 (colormap)", cmaps,
                            index=cmaps.index(ds.get("cmap", "viridis")) if ds.get("cmap", "viridis") in cmaps else 0, key="normal_cmap")

        spec_cols = st.columns(2)
        with spec_cols[0]:
            spec_upper_txt = st.text_input("上限 SPEC（可选，μm）", value="" if ds.get("spec_upper") is None else str(ds.get("spec_upper")), key="normal_spec_up_txt")
        with spec_cols[1]:
            spec_lower_txt = st.text_input("下限 SPEC（可选，μm）", value="" if ds.get("spec_lower") is None else str(ds.get("spec_lower")), key="normal_spec_lo_txt")

        def _parse(s):
            try:
                return float(s) if s.strip() else None
            except Exception:
                return None

        spec_upper = _parse(spec_upper_txt)
        spec_lower = _parse(spec_lower_txt)

        ds["wafer_inch"] = wafer_inch
        ds["cmap"] = cmap
        ds["spec_upper"] = spec_upper
        ds["spec_lower"] = spec_lower

        grid_res = st.slider("插值网格分辨率（越高越细，越慢）", 120, 360, 220, 10, key="normal_grid")
        show_labels = st.checkbox("顶视图显示厚度标签", value=True, key="normal_labels")

        with st.expander("顶视图视窗（用于导出 PPT 的“当前视图”）", expanded=False):
            cx = st.number_input("中心 X (mm)", value=0.0, step=1.0, key="tv_cx")
            cy = st.number_input("中心 Y (mm)", value=0.0, step=1.0, key="tv_cy")
            zoom = st.slider("缩放（越大越放大）", 0.5, 3.0, 1.0, 0.05, key="tv_zoom")

        df = ds["df"]

        st.markdown("#### 统计结果（μm）")
        stats_df = compute_statistics(df["thickness"].to_numpy(dtype=float))
        st.dataframe(stats_df, use_container_width=True, hide_index=True)

        st.markdown("---")
        st.markdown("#### 导出到一页 PPT（单击即下载）")
        ppt_title = st.text_input("PPT 标题", value=f"Wafer Thickness Report - {ds_name}", key="ppt_title")

        # NOTE: actual download_button is placed after figures are built on the right column (needs fig json),
        # but we keep the title input here.

    with right:
        df = ds["df"]

        st.markdown("#### 数据表（X, Y, Radius, Thickness）")
        disp_df = with_radius_columns(df)
        st.dataframe(
            disp_df.rename(columns={"x": "X (mm)", "y": "Y (mm)", "radius": "Radius (mm)", "thickness": "Thickness (μm)"}),
            use_container_width=True,
            hide_index=True,
        )

        top_fig = make_top_view_heatmap(
            df,
            radius_mm=radius_mm,
            cmap=cmap,
            grid_res=grid_res,
            show_labels=show_labels,
            view_center=(st.session_state.get("tv_cx", 0.0), st.session_state.get("tv_cy", 0.0)),
            view_zoom=st.session_state.get("tv_zoom", 1.0),
        )
        st.plotly_chart(top_fig, use_container_width=True, key="top_view_fig")

        st.markdown("---")
        st.markdown("### 3D 视图")
        cam_box = st.expander("视角与 3D 显示设置（可选）", expanded=False)
        with cam_box:
            eye_cols = st.columns(3)
            eye_x = eye_cols[0].slider("eye.x", -3.0, 3.0, 1.7, 0.1, key="eye_x")
            eye_y = eye_cols[1].slider("eye.y", -3.0, 3.0, 1.7, 0.1, key="eye_y")
            eye_z = eye_cols[2].slider("eye.z", 0.1, 5.0, 1.2, 0.1, key="eye_z")
            camera = (eye_x, eye_y, eye_z)

            z_scale = st.slider("Z 方向夸张倍数（仅影响显示）", 1.0, 500.0, 10.0, 1.0, key="z_scale")
            z_aspect = st.slider("纵向视觉比例（仅影响显示）", 0.2, 10.0, 2.5, 0.1, key="z_aspect")

        fig3d = make_3d_surface(
            df,
            radius_mm=radius_mm,
            cmap=cmap,
            grid_res=grid_res,
            spec_upper=spec_upper,
            spec_lower=spec_lower,
            camera_eye=camera,
            z_scale=z_scale,
            z_aspect=z_aspect,
        )
        st.plotly_chart(fig3d, use_container_width=True, key="surface3d")

        # --- One-click PPT download (build bytes, then download) ---
        st.markdown("---")
        st.markdown("### 导出 PPT（单击即下载）")

        # Prepare export tables
        raw_export = with_radius_columns(df).rename(columns={"x": "X(mm)", "y": "Y(mm)", "radius": "R(mm)", "thickness": "T(μm)"}).copy()
        # Ensure X/Y/R are 2 decimals as strings in PPT
        raw_export["X(mm)"] = raw_export["X(mm)"].map(lambda v: f"{float(v):.2f}")
        raw_export["Y(mm)"] = raw_export["Y(mm)"].map(lambda v: f"{float(v):.2f}")
        raw_export["R(mm)"] = raw_export["R(mm)"].map(lambda v: f"{float(v):.2f}")
        raw_export["T(μm)"] = raw_export["T(μm)"].map(lambda v: f"{float(v):.4f}")

        stats_export = compute_statistics(df["thickness"].to_numpy(dtype=float)).copy()

        subtitle = f"Dataset: {ds_name} | Wafer: {wafer_inch:.0f} inch | {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"

        # cache keys (strings) to enable st.cache_data
        raw_csv = raw_export.to_csv(index=False)
        stats_csv = stats_export.to_csv(index=False)
        top_json = top_fig.to_json()
        surf_json = fig3d.to_json()

        # Build bytes (cached). If kaleido missing, this will throw -> show error.
        ppt_bytes = None
        err = None
        try:
            ppt_bytes = build_ppt_cached(st.session_state.get("ppt_title", f"Wafer Thickness Report - {ds_name}"), subtitle, raw_csv, stats_csv, top_json, surf_json)
        except Exception as e:
            err = e

        if err is not None:
            st.error(f"PPT 构建失败：{err}")
            st.info("若提示 kaleido 缺失：请确认 requirements.txt 里有 kaleido==0.2.1 并清缓存重启。")
        else:
            st.download_button(
                "📤 生成 PPT 并下载",
                data=ppt_bytes,
                file_name=f"{ds_name}_report.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True,
            )


# ============================================================
# Main
# ============================================================

def main():
    st.title("晶圆厚度 3D 可视化（专业版）")
    tab_normal, tab_gen = st.tabs(["📊 普通模式 (Analysis)", "🛠️ 坐标生成器 (Generator)"])
    with tab_normal:
        normal_mode_ui()
    with tab_gen:
        blind_mode_ui()


if __name__ == "__main__":
    main()
